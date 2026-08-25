#!/usr/bin/env python3
"""Extract text/tables/images from prototype files in file/ directory."""
import os
import json
import shutil
from pathlib import Path
from collections import defaultdict

ROOT = Path('D:/VibeTest/bigsound')
FILE_DIR = ROOT / 'file'
OUT_DIR = ROOT / 'public' / 'images' / 'prototype'
OUT_DIR.mkdir(parents=True, exist_ok=True)
REPORT_FILE = ROOT / 'scripts' / 'prototype_extraction_report.json'

def extract_excel(path: Path):
    import pandas as pd
    xl = pd.ExcelFile(path)
    sheets = {}
    for name in xl.sheet_names:
        try:
            df = pd.read_excel(xl, sheet_name=name, header=None)
            # Drop completely empty rows/cols
            df = df.dropna(how='all').dropna(axis=1, how='all')
            sheets[name] = df.fillna('').astype(str).values.tolist()
        except Exception as e:
            sheets[name] = f'ERROR: {e}'
    return {'type': 'excel', 'sheets': sheets}

def extract_pdf_text(path: Path):
    try:
        import pdfplumber
        pages = []
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages[:50], 1):  # limit first 50 pages
                text = page.extract_text() or ''
                tables = page.extract_tables() or []
                pages.append({'page': i, 'text': text[:5000], 'tables': tables[:5]})
        return {'type': 'pdf', 'pages': pages}
    except Exception as e:
        return {'type': 'pdf', 'error': str(e)}

def extract_pptx(path: Path):
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    texts.append(shape.text)
            slides.append({'slide': i, 'text': '\n'.join(texts)})
        return {'type': 'pptx', 'slides': slides}
    except Exception as e:
        return {'type': 'pptx', 'error': str(e)}

def extract_pdf_images(path: Path, prefix: str):
    """Use pdfimages if available; fallback to pymupdf render."""
    from pypdf import PdfReader
    reader = PdfReader(path)
    img_info = []
    for i, page in enumerate(reader.pages[:20], 1):
        try:
            imgs = page.images
            for j, img in enumerate(imgs[:10]):
                ext = img.name.split('.')[-1] if '.' in img.name else 'png'
                out_name = f'{prefix}_p{i}_img{j}.{ext}'
                (OUT_DIR / out_name).write_bytes(img.data)
                img_info.append(out_name)
        except Exception as e:
            pass
    return img_info

def extract_pptx_images(path: Path, prefix: str):
    from pptx import Presentation
    prs = Presentation(path)
    img_info = []
    for i, slide in enumerate(prs.slides, 1):
        for j, shape in enumerate(slide.shapes):
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    img = shape.image
                    ext = img.ext
                    out_name = f'{prefix}_s{i}_img{j}.{ext}'
                    (OUT_DIR / out_name).write_bytes(img.blob)
                    img_info.append(out_name)
                except Exception as e:
                    pass
    return img_info

def main():
    report = {}
    for f in sorted(FILE_DIR.iterdir()):
        print(f'Processing {f.name} ...')
        if f.suffix.lower() in ('.xlsx', '.xls'):
            report[f.name] = extract_excel(f)
        elif f.suffix.lower() == '.pdf':
            report[f.name] = extract_pdf_text(f)
            prefix = f.stem.replace(' ', '_').replace('.', '_')
            imgs = extract_pdf_images(f, prefix)
            report[f.name]['extracted_images'] = imgs
        elif f.suffix.lower() in ('.pptx', '.ppt'):
            report[f.name] = extract_pptx(f)
            prefix = f.stem.replace(' ', '_').replace('.', '_')
            imgs = extract_pptx_images(f, prefix)
            report[f.name]['extracted_images'] = imgs
        else:
            report[f.name] = {'type': 'unknown'}

    REPORT_FILE.parent.mkdir(parents=True, exist_ok=True)
    REPORT_FILE.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding='utf-8')
    print(f'Report saved to {REPORT_FILE}')
    print(f'Images saved to {OUT_DIR}')

if __name__ == '__main__':
    main()
